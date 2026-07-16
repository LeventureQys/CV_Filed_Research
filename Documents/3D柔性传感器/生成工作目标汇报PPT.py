from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REPO_ROOT = ROOT.parents[1]
ASSET_DIR = ROOT / "assets"
PPT_ASSET_DIR = ASSET_DIR / "工作目标汇报"
PPT_ASSET_DIR.mkdir(exist_ok=True)
OUTPUT_PATH = ROOT / "柔性压感传感器软件建模工作目标汇报.pptx"
NOTES_PATH = ROOT / "柔性压感传感器软件建模工作目标汇报_讲述提纲.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x13, 0x25, 0x3D)
BLUE = RGBColor(0x24, 0x67, 0xB2)
CYAN = RGBColor(0x20, 0xA7, 0xB5)
GREEN = RGBColor(0x2C, 0x8C, 0x68)
ORANGE = RGBColor(0xD4, 0x87, 0x2C)
RED = RGBColor(0xC8, 0x4D, 0x5F)
PURPLE = RGBColor(0x79, 0x5B, 0xA8)
TEXT = RGBColor(0x20, 0x2A, 0x38)
MUTED = RGBColor(0x65, 0x72, 0x84)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xE1, 0xEA)

FONT = "Microsoft YaHei"


def prepare_logo():
    source = REPO_ROOT / "assets" / "logo.ico"
    target = PPT_ASSET_DIR / "logo.png"
    with Image.open(source) as image:
        image.convert("RGBA").save(target)
    return target


LOGO_PATH = prepare_logo()
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
speaker_notes = []


def add_text(slide, text, left, top, width, height, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(3)
    frame.margin_right = Pt(3)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_paragraphs(slide, lines, left, top, width, height, size=17, color=TEXT, bullet=False, spacing=8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(spacing)
        paragraph.level = 0
        if bullet:
            paragraph.text = f"• {line}"
    return box


def add_shape(slide, left, top, width, height, fill, radius=True, line_color=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, title, body, left, top, width, height, accent=BLUE, title_size=18, body_size=13):
    add_shape(slide, left, top, width, height, WHITE, line_color=LINE)
    add_shape(slide, left, top, 0.08, height, accent, radius=False)
    add_text(slide, title, left + 0.22, top + 0.13, width - 0.4, 0.46, title_size, NAVY, True)
    add_text(slide, body, left + 0.22, top + 0.61, width - 0.4, height - 0.72, body_size, MUTED)


def add_pill(slide, text, left, top, width, fill, color=WHITE):
    shape = add_shape(slide, left, top, width, 0.38, fill)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = FONT
    paragraph.font.size = Pt(11)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color


def add_picture_contain(slide, path, left, top, width, height):
    with Image.open(path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio > box_ratio:
        draw_width = width
        draw_height = width / image_ratio
        draw_left = left
        draw_top = top + (height - draw_height) / 2
    else:
        draw_height = height
        draw_width = height * image_ratio
        draw_left = left + (width - draw_width) / 2
        draw_top = top
    slide.shapes.add_picture(str(path), Inches(draw_left), Inches(draw_top), Inches(draw_width), Inches(draw_height))


def add_connector(slide, x1, y1, x2, y2, color=BLUE, width=2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True


def add_brand(slide, page):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_shape(slide, 0, 0, 13.333, 0.08, CYAN, radius=False)
    slide.shapes.add_picture(str(LOGO_PATH), Inches(11.92), Inches(0.20), Inches(0.48), Inches(0.48))
    add_text(slide, "模量科技 ModuTech", 10.0, 0.27, 1.85, 0.28, 10, MUTED, True, PP_ALIGN.RIGHT)
    add_shape(slide, 0.65, 7.04, 12.02, 0.012, LINE, radius=False)
    add_text(slide, "柔性压感传感器软件建模工作目标", 0.68, 7.08, 5.2, 0.22, 9, MUTED)
    add_text(slide, str(page), 12.0, 7.08, 0.55, 0.22, 9, MUTED, False, PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=""):
    add_text(slide, title, 0.72, 0.42, 10.8, 0.55, 27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.75, 1.02, 11.5, 0.36, 13, MUTED)


def new_slide(title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    page = len(prs.slides)
    add_brand(slide, page)
    add_title(slide, title, subtitle)
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes
    speaker_notes.append((page, title, notes))
    return slide


def add_metric(slide, value, label, left, top, width, color):
    add_shape(slide, left, top, width, 1.28, WHITE, line_color=LINE)
    add_text(slide, value, left, top + 0.12, width, 0.55, 27, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, left + 0.1, top + 0.73, width - 0.2, 0.32, 11, MUTED, False, PP_ALIGN.CENTER)


# 1. Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY
add_shape(slide, 0, 0, 13.333, 0.12, CYAN, radius=False)
slide.shapes.add_picture(str(LOGO_PATH), Inches(0.82), Inches(0.62), Inches(0.78), Inches(0.78))
add_text(slide, "模量科技 ModuTech", 1.72, 0.78, 3.2, 0.4, 17, WHITE, True)
add_text(slide, "柔性压感传感器\n软件建模工作目标", 0.82, 2.05, 7.4, 1.55, 34, WHITE, True)
add_text(slide, "不修改工艺条件下，上位机与嵌入式如何提高可解释性、稳定性与可标定性", 0.86, 3.78, 9.8, 0.55, 17, RGBColor(0xC5, 0xD5, 0xE8))
add_pill(slide, "向上汇报 / 下一步工作计划", 0.86, 4.62, 2.5, CYAN)
add_text(slide, "汇报人：李文轩", 0.86, 5.65, 3.0, 0.38, 15, WHITE, True)
add_text(slide, "面向：公司管理层与研发负责人", 0.86, 6.08, 4.6, 0.34, 12, RGBColor(0xC5, 0xD5, 0xE8))
add_shape(slide, 9.15, 1.55, 3.0, 3.9, RGBColor(0x1A, 0x35, 0x56))
add_text(slide, "目标", 9.55, 1.96, 2.2, 0.4, 17, CYAN, True, PP_ALIGN.CENTER)
add_text(slide, "把“数据不好看”\n转化为\n可分解的问题\n可验证的模型\n可落地的软件管线", 9.5, 2.58, 2.3, 2.2, 20, WHITE, True, PP_ALIGN.CENTER)
cover_notes = "开场先说明：这次不是汇报某一个滤波算法，而是汇报下一阶段如何把柔性压感中的时漂、噪声、跨次差异和串扰，转化为上位机与嵌入式可执行的软件建模工作。工艺暂不修改，重点是把能通过软件解决的问题做清楚，把不能解决的边界也讲清楚。"
slide.notes_slide.notes_text_frame.text = cover_notes
speaker_notes.append((1, "封面", cover_notes))


# 2. Executive summary
slide = new_slide(
    "一页结论：下一阶段不是“继续调滤波”",
    "建议建立一套按问题分层、按设备分流、可回放验收的软件建模体系",
    "先给老板结论：当前的问题至少包含五类，滤波只能解决其中一小部分。下一阶段的价值不是把曲线做得更平，而是建立可解释、可标定、可复现的测量链路。请求批准的工作分三阶段：链路透明、模型辨识、产品化补偿。",
)
add_card(slide, "当前困境", "时漂、帧间噪声、约 15% 跨次差异、空间伪影、正压非零背景", 0.72, 1.58, 3.75, 1.48, RED)
add_card(slide, "核心判断", "统一的应是建模层级，不是所有设备共用同一套参数", 4.79, 1.58, 3.75, 1.48, ORANGE)
add_card(slide, "软件目标", "保留原始量，分离误差来源，建立校正量、力值与质量状态", 8.86, 1.58, 3.75, 1.48, GREEN)
add_metric(slide, "P0", "链路透明与数据规范", 0.85, 3.55, 2.55, BLUE)
add_connector(slide, 3.48, 4.18, 4.14, 4.18, MUTED)
add_metric(slide, "P1", "方差分解与灰盒辨识", 4.22, 3.55, 2.55, ORANGE)
add_connector(slide, 6.85, 4.18, 7.51, 4.18, MUTED)
add_metric(slide, "P2", "证据充分后的高级补偿", 7.59, 3.55, 2.55, GREEN)
add_text(slide, "阶段门：每一步都以留出数据上的可重复收益为准，而不是算法复杂度。", 1.0, 5.45, 11.3, 0.48, 17, NAVY, True, PP_ALIGN.CENTER)


# 3. Problems overview
slide = new_slide(
    "我们面对的是五类数学问题",
    "这些问题的时间尺度、空间尺度和可修复性完全不同",
    "这一页建立共同语言。请注意跨次差异不等于随机噪声，时漂也不等于基线漂移。后面的研究和软件架构会按这五类问题分别设计，避免用一个滤波器处理所有现象。",
)
problems = [
    ("零点与背景", "空载非零、共模、工频", BLUE),
    ("静态映射", "非线性、饱和、Cell 增益", CYAN),
    ("记忆与时漂", "迟滞、蠕变、恢复、温漂", ORANGE),
    ("阵列与空间", "机械扩散、串扰、坏点", PURPLE),
    ("显示与交互", "跳动、闪烁、释放残影", GREEN),
]
for index, (title, body, color) in enumerate(problems):
    left = 0.78 + index * 2.48
    add_shape(slide, left, 1.75, 2.12, 2.15, WHITE, line_color=LINE)
    add_shape(slide, left, 1.75, 2.12, 0.12, color, radius=False)
    add_text(slide, f"0{index + 1}", left + 0.18, 2.02, 0.5, 0.4, 14, color, True)
    add_text(slide, title, left + 0.18, 2.46, 1.78, 0.45, 17, NAVY, True)
    add_text(slide, body, left + 0.18, 3.04, 1.78, 0.6, 12, MUTED)
add_shape(slide, 1.1, 4.65, 11.15, 1.12, RGBColor(0xE8, 0xF1, 0xFA), line_color=RGBColor(0xC3, 0xD7, 0xEA))
add_text(slide, "原则：先识别问题属于哪一类，再决定是标定、状态模型、空间校正、时间滤波，还是只输出质量标志。", 1.35, 4.85, 10.65, 0.72, 17, NAVY, True, PP_ALIGN.CENTER)


# 4. Evidence: drift and noise
slide = new_slide(
    "困境证据：时漂和帧间噪声必须分开",
    "同一条曲线中可以同时存在慢趋势和高频跳动，但二者需要不同模型",
    "左图强调固定载荷下仍有慢趋势，右图强调稳态局部窗口的帧间震荡。低通只对右边有效；左边需要区分材料蠕变、界面变化、温度和电路预热。汇报时不要承诺滤波能够解决全部时漂。",
)
add_picture_contain(slide, ASSET_DIR / "困境-数据时漂.png", 0.62, 1.45, 6.12, 4.75)
add_picture_contain(slide, ASSET_DIR / "困境-数据噪声.png", 6.78, 1.45, 5.93, 4.75)
add_pill(slide, "慢趋势 → 状态/环境模型", 1.6, 6.25, 2.4, ORANGE)
add_pill(slide, "帧间震荡 → 因果低通", 8.45, 6.25, 2.4, GREEN)


# 5. Evidence: crosstalk and repeatability
slide = new_slide(
    "困境证据：空间伪影与跨次差异影响“可信度”",
    "热图更干净不代表真实压力被恢复；同一砝码更平滑也不代表两次结果一致",
    "左侧是串扰机理示意，必须明确它不是当前膜片的实测系数。右侧解释为什么 15% 跨次差异不能靠 EMA 解决：它可能来自位置、装夹、材料历史、环境和设备。我们的任务是用实验把这些方差拆开。",
)
add_picture_contain(slide, ASSET_DIR / "困境-串扰伪影.png", 0.65, 1.52, 6.35, 4.65)
add_card(slide, "约 15% 跨次差异", "对触摸检测可能可容忍；对定量测力通常意味着尚未形成稳定标定。", 7.35, 1.72, 5.15, 1.42, RED, 20, 14)
add_card(slide, "不是同一种“噪声”", "保持段波动、重复放置、位置差异、跨日期和跨设备必须分别统计。", 7.35, 3.43, 5.15, 1.42, ORANGE, 18, 14)
add_card(slide, "软件可做什么", "通过元数据、受控实验、分层参数和质量标志，把随机差异转成可判断的工程状态。", 7.35, 5.14, 5.15, 1.42, GREEN, 18, 14)


# 6. Boundary
slide = new_slide(
    "不修改工艺，软件能解决多少？",
    "必须明确能力边界，避免把不可辨识问题包装成算法效果",
    "这一页用于管理预期。软件能够修复的是可观测、可重复、可验证的误差；对于随机接触、结构变化和已被截断的信息，只能通过质量标志、重新标定或推动硬件链路改进，不能承诺恢复真实力。",
)
add_card(slide, "可直接校正", "软件基线、已知偏置、可观测共模、板级 offset/gain、稳定非线性", 0.78, 1.7, 3.75, 2.15, GREEN, 20, 14)
add_card(slide, "可部分建模", "迟滞、蠕变、温漂、位置响应；前提是有状态、温度或实验变量", 4.79, 1.7, 3.75, 2.15, ORANGE, 20, 14)
add_card(slide, "不可可靠恢复", "随机装夹、不可重复材料跳变、模拟饱和、已截零且未保存的原始信息", 8.80, 1.7, 3.75, 2.15, RED, 20, 14)
add_shape(slide, 1.1, 4.45, 11.15, 1.3, NAVY)
add_text(slide, "软件工作的价值 = 提升可解释性 + 提升可标定性 + 减少可观测误差 + 明确不可用状态", 1.4, 4.72, 10.55, 0.76, 20, WHITE, True, PP_ALIGN.CENTER)


# 7. Modeling framework
slide = new_slide(
    "理论目标：建立分层灰盒模型",
    "统一的是观测链路，不是碳膜、织物和手套共用同一组参数",
    "从真实载荷到产品输出，依次经过结构传播、材料状态、电学响应、采集电路和软件校正。我们不追求一开始建立完整材料有限元，而是建立每个参数都能通过实验辨识的灰盒模型。",
)
add_picture_contain(slide, ASSET_DIR / "分层模型.png", 0.72, 1.4, 7.05, 5.35)
add_card(slide, "快 / 慢状态", "描述快速接触建立与慢速蠕变、恢复；加载、卸载、空载预热分别拟合。", 8.15, 1.62, 4.45, 1.35, ORANGE, 18, 13)
add_card(slide, "Cell 分层参数", "全局模型 + 设备修正 + Cell 修正，避免每个 Cell 独立过拟合。", 8.15, 3.25, 4.45, 1.35, BLUE, 18, 13)
add_card(slide, "空间观测模型", "先验证线性、叠加和位置稳定，再决定局部核、分区核或只做质量提示。", 8.15, 4.88, 4.45, 1.35, PURPLE, 18, 13)


# 8. Research method
slide = new_slide(
    "研究方法：用受控实验完成方差分解",
    "先回答误差来自哪里，再决定模型复杂度",
    "这页讲具体研究方法。每组实验只改变一个主要因素，并完整保留空载、加载、保持、释放和恢复。训练与测试按整条录制、位置、日期和设备划分，不能用测试录制自身估计最终参数。",
)
steps = [
    ("01", "链路审计", "确认偏置、整流、截零、滤波发生在哪一层", BLUE),
    ("02", "噪声分类", "固定电阻 / 空载 / 多次上电 / 不同采样率", CYAN),
    ("03", "动态辨识", "升降载、长保持、释放恢复、温度同步记录", ORANGE),
    ("04", "空间辨识", "中心边角、不同压头、单点扫描、双点叠加", PURPLE),
    ("05", "留出验收", "完整录制、位置、日期、设备独立测试", GREEN),
]
for index, (number, title, body, color) in enumerate(steps):
    left = 0.72 + index * 2.48
    add_shape(slide, left, 1.72, 2.13, 3.3, WHITE, line_color=LINE)
    add_text(slide, number, left + 0.18, 1.98, 0.55, 0.38, 14, color, True)
    add_text(slide, title, left + 0.18, 2.53, 1.75, 0.48, 18, NAVY, True)
    add_text(slide, body, left + 0.18, 3.15, 1.75, 1.25, 13, MUTED)
    if index < len(steps) - 1:
        add_connector(slide, left + 2.13, 3.38, left + 2.45, 3.38, MUTED, 1.5)
add_text(slide, "输出：基线、增益、非线性、迟滞、快慢时间常数、温漂、空间响应、坏点与不确定度", 0.95, 5.62, 11.4, 0.55, 16, NAVY, True, PP_ALIGN.CENTER)


# 9. Embedded vs host
slide = new_slide(
    "软件分工：嵌入式保底，上位机建模",
    "低延迟确定性工作下沉；复杂、可配置、可回放模型留在上位机",
    "嵌入式重点保证数据可信和设备健康，上位机负责设备模式、模型选择和较重计算。两端都必须保留原始流，所有不可逆处理可以关闭。这样既满足实时性，也支持后续模型迭代。",
)
add_shape(slide, 0.78, 1.62, 5.65, 4.62, WHITE, line_color=LINE)
add_shape(slide, 6.88, 1.62, 5.65, 4.62, WHITE, line_color=LINE)
add_pill(slide, "嵌入式软件", 2.18, 1.92, 2.8, BLUE)
add_pill(slide, "上位机软件", 8.28, 1.92, 2.8, GREEN)
add_paragraphs(slide, [
    "时间戳、帧序号、丢帧与采样率",
    "ADC offset/gain、通道映射、坏点表",
    "饱和、断路、短路、卡死与参考监测",
    "可关闭的轻量去毛刺、陷波和 EMA",
    "原始有符号数据与质量标志持续上报",
], 1.18, 2.62, 4.85, 2.9, 14, TEXT, True, 9)
add_paragraphs(slide, [
    "设备 / 采集模式识别与参数版本管理",
    "基线、共模、坏点和静态非线性标定",
    "快慢状态、迟滞、温漂与总力估计",
    "有证据时执行空间响应校正",
    "测量流、记录流和显示流彻底分离",
], 7.28, 2.62, 4.85, 2.9, 14, TEXT, True, 9)
add_text(slide, "共同原则：原始数据可获取 · 参数可追溯 · 算法可回放 · 处理可绕过", 1.2, 5.72, 10.95, 0.4, 15, NAVY, True, PP_ALIGN.CENTER)


# 10. Host pipeline
slide = new_slide(
    "目标软件管线：从原始帧到可解释产品输出",
    "每一步都有输入、参数、质量标志和可独立验收的指标",
    "这页是上位机开发主线。先校验和识别设备，再做基线与坏点，之后才进入空间、标定和动态状态。显示处理必须最后分支，不能反写测量数据。每个模块都应支持离线回放和 A/B。",
)
add_picture_contain(slide, ASSET_DIR / "上位机核心算法.png", 0.7, 1.38, 12.0, 5.35)


# 11. Plan
slide = new_slide(
    "下一步工作计划：三阶段、三个阶段门",
    "以证据推进，而不是以算法复杂度推进",
    "这是本次汇报希望确认的工作计划。P0 解决数据和链路透明；P1 完成模型辨识，是主要工作量；P2 只在证据充分时做高级补偿。每个阶段都有明确产出和继续投入的条件。",
)
plans = [
    ("P0", "链路透明", "2–3 周", "观测链路表\n原始数据规范\n统一回放基线", BLUE),
    ("P1", "模型辨识", "4–6 周", "方差分解\n逐 Cell 标定\n动态 / 空间模型", ORANGE),
    ("P2", "产品化补偿", "证据门后", "设备分流配置\n在线状态模型\n跨设备验证", GREEN),
]
for index, (phase, title, duration, body, color) in enumerate(plans):
    left = 0.82 + index * 4.18
    add_shape(slide, left, 1.72, 3.72, 3.8, WHITE, line_color=LINE)
    add_pill(slide, phase, left + 0.22, 1.98, 0.82, color)
    add_text(slide, duration, left + 2.0, 2.02, 1.35, 0.3, 12, color, True, PP_ALIGN.RIGHT)
    add_text(slide, title, left + 0.24, 2.62, 3.1, 0.5, 22, NAVY, True)
    add_text(slide, body, left + 0.24, 3.36, 3.1, 1.15, 15, MUTED)
    if index < len(plans) - 1:
        add_connector(slide, left + 3.72, 3.58, left + 4.08, 3.58, MUTED, 2)
add_shape(slide, 1.0, 5.86, 11.3, 0.72, RGBColor(0xE8, 0xF1, 0xFA), line_color=RGBColor(0xC3, 0xD7, 0xEA))
add_text(slide, "阶段门：参数可重复、留出集有效、跨设备收益稳定，才进入下一阶段。", 1.3, 6.06, 10.7, 0.34, 16, NAVY, True, PP_ALIGN.CENTER)


# 12. Outcomes and ask
slide = new_slide(
    "最终效果与本次决策请求",
    "目标不是承诺“消灭材料差异”，而是把系统从不可解释推进到可测量、可诊断、可迭代",
    "最后收口到管理价值和决策请求。最终应能明确哪些误差已经补偿、哪些状态需要重新标定、哪些设备不可用。请求确认产品目标、允许固件保留原始数据、协调硬件提供链路信息，并批准受控实验资源。",
)
outcomes = [
    ("数据可信", "原始、校正、力值和质量状态可追溯", BLUE),
    ("模型可解释", "能区分噪声、漂移、迟滞、位置和设备差异", ORANGE),
    ("产品可配置", "碳膜、织物、手套按设备和模式分流", GREEN),
    ("效果可验收", "准确度、稳定性、延迟和空间指标同时报告", PURPLE),
]
for index, (title, body, color) in enumerate(outcomes):
    left = 0.78 + (index % 2) * 6.05
    top = 1.62 + (index // 2) * 1.58
    add_card(slide, title, body, left, top, 5.62, 1.28, color, 18, 13)
add_shape(slide, 0.9, 4.92, 11.55, 1.22, NAVY)
add_text(slide, "需要确认", 1.25, 5.17, 1.25, 0.38, 15, CYAN, True)
add_text(slide, "① 产品目标与验收指标   ② 原始数据链路支持   ③ 受控实验与跨设备样本   ④ P0 / P1 工作优先级", 2.45, 5.05, 9.45, 0.72, 15, WHITE, True)
add_text(slide, "李文轩｜模量科技 ModuTech", 0.9, 6.48, 4.0, 0.3, 11, MUTED, True)


prs.save(OUTPUT_PATH)

notes_lines = [
    "# 柔性压感传感器软件建模工作目标汇报｜讲述提纲",
    "",
    "> 建议时长：15–20 分钟。每页先讲结论，再讲图中的一项证据，最后自然过渡到下一页。",
    "",
]
for page, title, notes in speaker_notes:
    notes_lines.extend([f"## {page}. {title}", "", notes, ""])
NOTES_PATH.write_text("\n".join(notes_lines), encoding="utf-8")

print(f"PPT saved: {OUTPUT_PATH}")
print(f"Notes saved: {NOTES_PATH}")
print(f"Slides: {len(prs.slides)}")
