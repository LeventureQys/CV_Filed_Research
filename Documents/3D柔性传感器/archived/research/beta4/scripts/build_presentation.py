from __future__ import annotations

import os
import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "output" / "ppt_assets"
OUTPUT = ROOT / "output"
RESEARCH_DIR = ROOT.parent
sys.path.insert(0, str(ROOT / "scripts"))
sys.path.insert(0, str(RESEARCH_DIR / "src"))

from data_loader_v2 import load_recording
from filters import ema
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLUE_DARK = RGBColor(0x1A, 0x3C, 0x6E)
BLUE_ACCENT = RGBColor(0x2E, 0x86, 0xDE)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def _build_background_comparison_chart() -> Path:
    chart_path = ASSETS / "background_raw_vs_ema.png"
    zoom_chart_path = ASSETS / "background_raw_vs_ema_zoom.png"
    background_path = RESEARCH_DIR / "DataSet" / "织物垫" / "背景噪声" / "录制数据_20260710143202_part0.csv"
    stress_path = RESEARCH_DIR / "DataSet" / "织物垫" / "1kg砝码压力" / "录制数据_20260710143635_part0.csv"

    background, _ = load_recording(str(background_path))
    stress, metadata = load_recording(str(stress_path))
    baseline = np.median(background, axis=0)
    corrected = np.maximum(stress - baseline, 0.0)
    raw = corrected.sum(axis=1)
    sample_rate = float(metadata["sample_freq_hz"])
    optimized = ema(raw, sample_rate, 0.8)
    time = np.arange(len(raw), dtype=float) / sample_rate

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    figure, axis = plt.subplots(figsize=(12.0, 2.5), constrained_layout=True)
    axis.plot(time, raw, color="#8A8A8A", linewidth=1.0, alpha=0.75, label="原始数据")
    axis.plot(time, optimized, color="#F28E2B", linewidth=2.2, label="EMA 0.8 s 优化后")
    axis.set_xlabel("时间（秒）")
    axis.set_ylabel("背景扣除后总 ADC")
    axis.set_title("1 kg 稳态实测：原始抖动与优化结果")
    axis.grid(alpha=0.2)
    axis.legend(loc="upper left", frameon=False)
    figure.savefig(chart_path, dpi=180, facecolor="white")
    plt.close(figure)

    stable_start_s = 12.0
    stable_end_s = 20.0
    stable_mask = (time >= stable_start_s) & (time <= stable_end_s)
    zoom_time = time[stable_mask]
    zoom_raw = raw[stable_mask]
    zoom_optimized = optimized[stable_mask]

    figure, axis = plt.subplots(figsize=(12.0, 2.5), constrained_layout=True)
    axis.plot(zoom_time, zoom_raw, color="#8A8A8A", linewidth=1.2, alpha=0.8, label="原始数据")
    axis.plot(zoom_time, zoom_optimized, color="#F28E2B", linewidth=2.4, label="EMA 0.8 s 优化后")
    axis.set_xlabel("时间（秒）")
    axis.set_ylabel("总 ADC")
    axis.set_title("稳态局部放大：相邻帧抖动明显降低")
    axis.grid(alpha=0.2)
    axis.legend(loc="upper left", frameon=False, ncol=2)
    figure.savefig(zoom_chart_path, dpi=180, facecolor="white")
    plt.close(figure)
    return chart_path


_build_background_comparison_chart()


def _add_logo(slide):
    logo_path = ASSETS / "logo.png"
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(12.5), Inches(0.25), Inches(0.6), Inches(0.6))


def _add_footer(slide, page_num: int):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"TactileSense | beta4 织物垫实物研究    {page_num}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT
    # divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.95), Inches(12.333), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()


def _add_title(slide, title: str, subtitle: str = ""):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    if subtitle:
        txSub = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.0), Inches(0.5))
        tf2 = txSub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def _add_body(slide, text: str, left=0.8, top=1.7, width=11.5, height=4.8, font_size=16):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)


def _add_card(slide, text: str, left, top, width, height, fill_color, text_color=WHITE, font_size=18):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True


def _add_image(slide, img_rel_path: str, left, top, width, height=None):
    full = Path(ROOT) / img_rel_path
    if full.exists():
        kwargs = {"left": Inches(left), "top": Inches(top), "width": Inches(width)}
        if height is not None:
            kwargs["height"] = Inches(height)
        slide.shapes.add_picture(str(full), **kwargs)


def _add_arrow(slide, left, top, width, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_standalone_bullet(slide, text, left, top, font_size=14, bold=False, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width := 11.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color


def _new_slide():
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _add_logo(slide)
    return slide


def _finalize(slide, num):
    _add_footer(slide, num)


# ============================================================
# Page 1: Cover
# ============================================================
s = _new_slide()
_add_title(s, "织物垫力值稳定算法研究", "TactileSense beta4 — 实物数据验证与算法选型")
_add_body(s, "汇报日期：2026-07-14\n基于 7 条织物垫实物录制，覆盖 500 g / 1 kg / 1.5 kg 及长时保持场景")
_finalize(s, 1)

# ============================================================
# Page 2: 一页结论
# ============================================================
s = _new_slide()
_add_title(s, "核心结论", "推荐 EMA 0.8 s 为默认算法")
_add_card(s, "推荐默认\nEMA 0.8 s", 0.8, 2.0, 2.8, 1.5, GREEN)
_add_card(s, "备选\nIIR 0.5 Hz / Huber EMA", 4.2, 2.0, 2.8, 1.5, BLUE_ACCENT)
_add_card(s, "预研中\n双状态卡尔曼", 7.6, 2.0, 2.8, 1.5, RGBColor(0xF3, 0x9C, 0x12))
_add_card(s, "已淘汰\n自适应基线", 10.8, 2.0, 2.0, 1.5, RED)
_add_body(s, "• 全部 12 组算法在 7 条录制上完成实物验证\n• 4 组入围基线推荐，其余已在指标表中说明放弃理由\n• 完整指标 CSV 与精度曲线见输出目录", top=4.0, font_size=15)
_finalize(s, 2)

# ============================================================
# Page 3: 问题背景
# ============================================================
s = _new_slide()
_add_title(s, "问题背景", "稳定受压时总力仍存在可见跳动")
_add_body(s, """• 织物垫传感器在放置砝码后，总 ADC 并非恒定不变
• 逐帧存在明显的数值跳动，影响显示稳定性和读数可信度
• 用户端感知为“数值不稳定”，即使施加的砝码质量未变化

本阶段目标：在砝码加载的稳态段，尽可能减少总力的显示跳动，
同时不显著改变稳态均值（即不“吃掉”真实力值信号）。""",
          left=0.8, top=1.75, width=11.4, height=4.7, font_size=17)
_finalize(s, 3)

# ============================================================
# Page 4: 实测优化效果
# ============================================================
s = _new_slide()
_add_title(s, "实测优化效果", "完整趋势与稳态局部放大采用同一条 1 kg 实物录制")
_add_image(s, "output/ppt_assets/background_raw_vs_ema.png", 1.0, 1.35, 11.3)
_add_standalone_bullet(s, "完整趋势：橙线贯穿全程并保持原始变化趋势", 1.0, 3.75, 12, True, BLUE_DARK)
_add_image(s, "output/ppt_assets/background_raw_vs_ema_zoom.png", 1.0, 4.05, 11.3)
_add_standalone_bullet(s, "稳态放大：灰线逐帧抖动明显，EMA 0.8 s 将其压低且均值基本不变", 1.0, 6.45, 12, True, BLUE_DARK)
_finalize(s, 4)

# ============================================================
# Page 5: 问题剖析
# ============================================================
s = _new_slide()
_add_title(s, "问题剖析", "跳动不能混为一谈，需分离四种物理来源")
_add_card(s, "随机噪声\nADC 量化+独立 Cell 噪声", 0.8, 1.8, 5.5, 1.3, BLUE_DARK, WHITE, 14)
_add_card(s, "共同偏置\n固定背景 + Cell 间共同漂移", 6.7, 1.8, 5.5, 1.3, BLUE_ACCENT, WHITE, 14)
_add_card(s, "慢漂移\n温度、供电、蠕变引起的缓慢变化", 0.8, 3.5, 5.5, 1.3, RGBColor(0xF3, 0x9C, 0x12), WHITE, 14)
_add_card(s, "偶发尖峰\n电磁干扰、接触不良", 6.7, 3.5, 5.5, 1.3, RED, WHITE, 14)
_add_body(s, "• 时间平滑噪声最有效，但无法消除偏置和慢漂移\n• 不同来源需要不同对策：基线扣除、低通滤波、漂移建模、离群抑制\n• 本阶段聚焦偏置正确扣除后的随机噪声压降", top=5.2, font_size=14)
_finalize(s, 5)

# ============================================================
# Page 6: 研究设计
# ============================================================
s = _new_slide()
_add_title(s, "研究设计", "3 类场景 × 7 条录制 × 12 组算法")
_add_body(s, """数据来源：织物垫 96 Cell（8×12）传感器，ADC 模式

三类场景：
  ① 旧格式 (10 Hz) — 500 g / 1 kg 连续压力录制
  ② 新格式 (55 Hz) — 500 g / 1 kg / 1.5 kg 空载+加载+恢复循环
  ③ 长时 / 空载 — 1 kg 20 min 保持 + 空载 20 min

12 组算法：原始 / EMA 0.35s / EMA 0.8s / M3+EMA / IIR1/IIR2 /
  卡尔曼快/稳 / 逐Cell M3+EMA / 自适应基线 / 双状态卡尔曼 / Huber EMA""", top=1.7, font_size=15)
_finalize(s, 6)

# ============================================================
# Page 7: 横向算法调研
# ============================================================
s = _new_slide()
_add_title(s, "横向算法调研与筛选", "从简单到复杂，评估三个价值维度")
_add_body(s, """评估维度：① roughness 降低率 ② 均值保持率 ③ 实现复杂度

─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─ ─
  EMA / IIR： 一阶/二阶低通，参数秒/Hz 定义，实现最简单
  Median-3 + EMA： 非线性离群抑制，然后低通
  标量卡尔曼： 单状态随机游走，等价于自适应系数 EMA
  双状态卡尔曼： 总力 + 漂移两状态，建模更精细
  自适应基线： 滑动窗口逐 Cell 中位数扣背景
  Huber EMA： 对离群残差做截断的标准 EMA""", top=1.7, font_size=14)
_finalize(s, 7)

# ============================================================
# Page 8: 第一轮淘汰
# ============================================================
s = _new_slide()
_add_title(s, "第一轮淘汰", "哪些算法不适合织物垫砝码加载场景")
_add_card(s, "自适应基线\n均值损失 > 70%", 0.8, 1.8, 3.6, 1.3, RED, WHITE, 13)
_add_card(s, "M3 + EMA\n与 EMA 等价", 4.9, 1.8, 3.6, 1.3, GRAY, WHITE, 13)
_add_card(s, "标量卡尔曼\n不比 IIR 更好", 9.0, 1.8, 3.6, 1.3, GRAY, WHITE, 13)
_add_body(s, """• 自适应基线：织物垫受力 Cell 占比高（30-50%），滑动中位数混入信号，均值从 17378 降至 576，不适用
• Median-3 + EMA（总力级 & 逐Cell）：96 Cell 求和后独立离群已被自然平滑，M3 无额外收益
• 标量卡尔曼（快/稳态型）：与 EMA 0.8 s 表现接近但复杂度更高，beta3 已说明无额外优势""", top=3.6, font_size=14)
_finalize(s, 8)

# ============================================================
# Page 9: 效果展示
# ============================================================
s = _new_slide()
_add_title(s, "效果展示", "3 个代表场景：500 g / 1 kg / 1.5 kg 空载-加载-恢复循环")
# 已有 fabric_final_comparison.png 显示 7 条录制 × 4 组算法
# 对于管理层汇报，原始图包含子图多，截取前 3 个场景（新格式 500g / 1kg / 1.5kg）
_add_image(s, "output/fabric_final_comparison.png", 0.5, 1.6, 12.3, 5.5)
_finalize(s, 9)

# ============================================================
# Page 10: 指标展示
# ============================================================
s = _new_slide()
_add_title(s, "算法指标汇总", "分场景 roughness 降低率与均值保持率")
_add_image(s, "output/fabric_metrics_summary.png", 0.5, 1.6, 12.3, 5.5)
_finalize(s, 10)

# ============================================================
# Page 11: 决策建议
# ============================================================
s = _new_slide()
_add_title(s, "决策建议", "四层结论")
_add_card(s, "默认\nEMA 0.8 s", 0.8, 1.8, 2.5, 1.3, GREEN, WHITE, 16)
_add_card(s, "备选\nIIR / Huber", 3.8, 1.8, 2.5, 1.3, BLUE_ACCENT, WHITE, 16)
_add_card(s, "预研\n双状态卡尔曼", 6.8, 1.8, 2.5, 1.3, RGBColor(0xF3, 0x9C, 0x12), WHITE, 16)
_add_card(s, "淘汰\n自适应基线", 9.8, 1.8, 2.5, 1.3, RED, WHITE, 16)
_add_body(s, """EMA 0.8 s
  • roughness 降低 70-80%，均值保持 > 99%
  • 实现最简单，参数含义清晰
  • 推荐作为产品化默认方案

二阶 IIR 0.5 Hz / Huber EMA
  • 效果与 EMA 0.8 s 接近，可互换
  • Huber 在尖峰环境中提供额外保护
  • 作为备选保留，无需同时上线

双状态卡尔曼
  • roughness 最低，但均值高 1-5%
  • 需要更多数据校准 q_drift
  • 建议作为线下一轮预研方向""", top=3.5, font_size=14)

# 决策矩阵小卡片：场景适配
txBox = s.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "关键判断：选默认方案时不追最低 roughness，优先保证均值不变——EMA 0.8 s 在全部 7 条录制上均值偏差 < 1%。"
p.font.size = Pt(12)
p.font.color.rgb = BLUE_DARK
p.font.bold = True
_finalize(s, 11)

# ============================================================
# Page 12: 下一步
# ============================================================
s = _new_slide()
_add_title(s, "下一步", "产品化验证与数据采集规范")
_add_body(s, """1. 在 TactileSense 主程序中实现可配置时间常数的 EMA 滤波器
2. 每条录制确保自含空载段（本研究发现独立背景扣除在不同批次存在偏差）
3. 双状态卡尔曼在固定砝码数据上校准 q_drift，目标均值偏移 < 1%
4. 以同样管线对 64×64 膜片和手套进行实物验证
5. 本阶段未研究快速敲击、滑动——这些场景需另行评估""", top=1.7, font_size=16)
_finalize(s, 12)

# ============================================================
output_path = ROOT / "织物垫实物研究汇报.pptx"
try:
    prs.save(str(output_path))
except PermissionError:
    output_path = ROOT / "织物垫实物研究汇报_更新版.pptx"
    try:
        prs.save(str(output_path))
    except PermissionError:
        output_path = ROOT / "织物垫实物研究汇报_更新版2.pptx"
        prs.save(str(output_path))
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
